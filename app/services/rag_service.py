import os
from typing import List, Dict, Any
from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document

class RagService:
    def __init__(self):
        # Initialize a fast, local embedding model
        # This will download the model weights (approx ~90MB) on the first run
        self.embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        # Using an in-memory chroma DB for this session
        self.vector_store = Chroma(embedding_function=self.embeddings)

    def process_text(self, text: str) -> int:
        """Processes raw text, chunks it, and adds it to the vector store."""
        docs = [Document(page_content=text, metadata={"source": "user_input"})]
        return self._add_documents(docs)

    def process_file(self, file_path: str) -> int:
        """Processes a file (TXT, PDF, DOCX, PPTX), chunks it, and adds to vector store."""
        ext = file_path.lower().split('.')[-1]
        
        if ext == 'pdf':
            loader = PyPDFLoader(file_path)
            docs = loader.load()
        elif ext in ['txt', 'csv']:
            loader = TextLoader(file_path, encoding='utf-8')
            docs = loader.load()
        elif ext == 'docx':
            import docx
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            docs = [Document(page_content=text, metadata={"source": file_path})]
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text_lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_lines.append(shape.text)
            docs = [Document(page_content="\n".join(text_lines), metadata={"source": file_path})]
        else:
            raise ValueError(f"Unsupported file type for RAG: {file_path}")

        return self._add_documents(docs)

    def _add_documents(self, docs: List[Document]) -> int:
        """Helper to split and store documents."""
        chunks = self.text_splitter.split_documents(docs)
        if not chunks:
            return 0
        
        # Clear existing collection to start fresh for a new upload (optional, but good for this demo)
        self.vector_store = Chroma(embedding_function=self.embeddings)
        self.vector_store.add_documents(chunks)
        return len(chunks)

    def retrieve_context(self, query: str, k: int = 3) -> str:
        """Retrieves top-k relevant chunks for a given query."""
        if not self.vector_store:
            return ""
        
        # We can use similarity_search
        results = self.vector_store.similarity_search(query, k=k)
        
        # Combine the chunks into a single context string
        context = "\n\n".join([f"--- {i+1} ---\n{doc.page_content}" for i, doc in enumerate(results)])
        return context

# Instantiate a global service instance
rag_service = RagService()
