import json
import os
import shutil
import tempfile
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Depends
from pydantic import BaseModel
from app.models.schemas import ChatRequest, ChatResponse
from app.services.data_service import DataService
from app.services.ai_service import AIService
from app.services.rag_service import rag_service
from app.services.media_service import media_service

router = APIRouter()

# Global state
app_state = {
    "df": None,
    "summary": None,
    "is_rag_active": False,
}

class TextUploadRequest(BaseModel):
    text: str

@router.post("/upload/unified")
async def upload_unified(file: UploadFile = File(...)):
    """Unified upload endpoint for any file type."""
    filename = file.filename or ""
    ext = filename.lower().split('.')[-1]
    
    # Save the file temporarily
    fd, temp_file_path = tempfile.mkstemp(suffix=f".{ext}")
    os.close(fd)
    
    try:
        with open(temp_file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
            
        report_markdown = ""
        report_charts = []
        
        # Structured Data
        if ext in ['csv', 'xlsx']:
            await file.seek(0) # Reset file pointer after shutil.copyfileobj
            df = await DataService.process_upload(file)
            app_state["df"] = df
            summary = DataService.generate_statistical_summary(df)
            app_state["summary"] = summary
            app_state["is_rag_active"] = False
            
            summary_json = json.dumps(summary)
            report_markdown = AIService.generate_unified_report('structured', summary_json)
            
            # Proactively generate charts for all suitable columns
            
            # 1. Categorical / Text columns
            categorical_cols = df.select_dtypes(include=['object', 'string', 'category']).columns
            for col in categorical_cols:
                # Limit to top 15 categories to keep charts readable
                counts = df[col].value_counts().head(15)
                if len(counts) > 1:
                    chart_type = "pie" if len(counts) <= 5 else "bar"
                    report_charts.append({
                        "col": col,
                        "type": chart_type,
                        "labels": counts.index.astype(str).tolist(),
                        "values": counts.values.tolist(),
                        "title": f"Distribution of {col} (Top {len(counts)})"
                    })
                    
            # 2. Numerical columns (discrete/small cardinality)
            numeric_cols = df.select_dtypes(include=['number']).columns
            for col in numeric_cols:
                if df[col].nunique() <= 15:
                    counts = df[col].value_counts().sort_index()
                    if len(counts) > 1:
                        report_charts.append({
                            "col": col,
                            "type": "bar",
                            "labels": counts.index.astype(str).tolist(),
                            "values": counts.values.tolist(),
                            "title": f"Distribution of {col}"
                        })
            
        # Media (Video/Audio)
        elif ext in ['mp4', 'mov', 'wav', 'mp3']:
            transcription = media_service.process_video_or_audio(temp_file_path)
            rag_service.process_text(transcription)
            
            app_state["is_rag_active"] = True
            app_state["df"] = None
            app_state["summary"] = None
            
            report_markdown = AIService.generate_unified_report('unstructured', transcription)
            
        # Documents (PDF, DOCX, PPTX, TXT)
        elif ext in ['pdf', 'docx', 'pptx', 'txt']:
            chunks_count = rag_service.process_file(temp_file_path)
            app_state["is_rag_active"] = True
            app_state["df"] = None
            app_state["summary"] = None
            
            # Retrieve some context for the summary report
            context = rag_service.retrieve_context("Give me an overview of this document", k=5)
            report_markdown = AIService.generate_unified_report('unstructured', context)
            
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")
            
        return {"message": "Success", "report_markdown": report_markdown, "report_charts": report_charts}
        
    except Exception as e:
        print(f"Unified upload failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)

@router.post("/upload/unified/text")
async def upload_unified_text(request: TextUploadRequest):
    """Unified upload for raw text."""
    try:
        rag_service.process_text(request.text)
        app_state["is_rag_active"] = True
        app_state["df"] = None
        app_state["summary"] = None
        
        report_markdown = AIService.generate_unified_report('unstructured', request.text)
        return {"message": "Success", "report_markdown": report_markdown}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/chat", response_model=ChatResponse)
async def chat_with_data(request: ChatRequest):
    """Unified chat endpoint."""
    try:
        if app_state["is_rag_active"]:
            # RAG Mode
            context = rag_service.retrieve_context(request.message, k=3)
            if not context:
                return ChatResponse(reply="I don't have enough context to answer that. Please upload more data.")
            
            reply = AIService.chat_with_data(request.message, context, is_rag=True)
            return ChatResponse(reply=reply)
        else:
            # Data Dictionary Mode
            if app_state["summary"] is None:
                raise HTTPException(status_code=400, detail="Please upload a dataset first.")
            
            summary_str = json.dumps(app_state["summary"])
            reply_dict = AIService.chat_with_data_and_analytics(
                request.message, summary_str, app_state["df"]
            )
            return ChatResponse(reply=reply_dict["reply"], chart_data=reply_dict.get("chart_data"))
    except Exception as e:
        print(f"Chat failed: {e}")
        raise HTTPException(status_code=500, detail="An error occurred while processing your request.")

class ExportRequest(BaseModel):
    markdown: str

@router.post("/export/docx")
async def export_docx(request: ExportRequest):
    """Exports markdown to DOCX format."""
    import docx
    import io
    from fastapi import Response
    
    doc = docx.Document()
    doc.add_heading('Analytics Report', 0)
    
    for line in request.markdown.split('\n'):
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.strip() == '':
            continue
        else:
            # Strip simple markdown like bold
            clean_line = line.replace('**', '')
            doc.add_paragraph(clean_line)
            
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    
    headers = {'Content-Disposition': 'attachment; filename="Data_Analytics_Report.docx"'}
    return Response(
        content=file_stream.read(),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers=headers
    )

@router.post("/export/pptx")
async def export_pptx(request: ExportRequest):
    """Exports markdown to PPTX format."""
    from pptx import Presentation
    import io
    from fastapi import Response
    
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Analytics Report"
    subtitle.text = "Generated by IDD Agent"
    
    bullet_slide_layout = prs.slide_layouts[1]
    current_slide = None
    tf = None
    
    for line in request.markdown.split('\n'):
        if line.startswith('# ') or line.startswith('## '):
            current_slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = current_slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = line.replace('#', '').strip()
            tf = body_shape.text_frame
            tf.text = "" # Clear default text
        elif current_slide and line.strip() != '':
            p = tf.add_paragraph()
            p.text = line.replace('**', '').replace('*', '').strip()
            
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    headers = {'Content-Disposition': 'attachment; filename="Data_Analytics_Report.pptx"'}
    return Response(
        content=file_stream.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers
    )
